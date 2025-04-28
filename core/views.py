from django.conf import settings
from django.shortcuts import render, get_object_or_404, redirect
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from django.shortcuts import render, redirect, get_object_or_404
from .models import Terminal,Document as DocModel, Printer
import PyPDF2  
import os
import zipfile
import xml.dom.minidom
from django.http import JsonResponse
from django.core.files.storage import FileSystemStorage
import pandas as pd


from django.shortcuts import render, get_object_or_404, redirect
from .models import TemplateDocument, PrintJob
from django.contrib.auth.decorators import login_required
import requests

# def upload_and_read_files(request,pk):
#     pages2 = None
#     terminal2 = get_object_or_404(Terminal, id=pk)  # Получаем терминал
#     price_per_page = terminal2.price_per_page_bw  # Цена за 1 страницу (пример)
#     if request.method == 'POST' and request.FILES.get('file'):
#         uploaded_file = request.FILES['file']
#         file = uploaded_file
#         is_color = request.POST.get("color_print") == "on"
        
#         fs = FileSystemStorage()
#         filename = fs.save(uploaded_file.name, uploaded_file)
#         file_path = fs.path(filename)

#         # Определение типа файла
#         file_extension = os.path.splitext(filename)[1].lower()

#         if file_extension == '.docx':
#             content = read_docx(file_path)
#             pages2 = get_word_page_count(file_path)
#         elif file_extension in ['.xls', '.xlsx']:
#             content = read_excel(file_path)
#         elif file_extension == '.txt':
#             content = read_txt(file_path)
#         else:
#             content = "Unsupported file format"
#         total_price = pages2 * price_per_page
 
#         print(type(pages2),pages2)
        
#         document = DocModel.objects.create(
#             title=file.name,
#             file=file,
#             pages=pages2,
#             terminal=terminal2,
#             is_color=is_color,
#             price_document=total_price
#         )
#         #  Перенаправляем на оплату, передавая ID документа
#         return redirect("payment", doc_id=document.id)
#         # return JsonResponse({'content': content, 'pages': pages})

#     return render(request, 'upload.html',{"terminal": terminal2})
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image
import pandas as pd

def upload_and_read_files(request, pk):
    terminal = get_object_or_404(Terminal, id=pk)
    if request.method == 'POST' and request.FILES.get('file'):
        uploaded_file = request.FILES['file']
        fs = FileSystemStorage()
        filename = fs.save(uploaded_file.name, uploaded_file)
        file_path = fs.path(filename)
        file_extension = os.path.splitext(filename)[1].lower()

        # Получаем параметры печати из формы
        is_color = request.POST.get("color_print") == "on"
        is_duplex = request.POST.get("duplex_print") == "on"
        
        # Определяем количество страниц для разных форматов
        pages = 1  # По умолчанию 1 страница
        
        try:
            if file_extension == '.docx':
                content = read_docx(file_path)
                pages = get_word_page_count(file_path)
                # pages = len(DocxDocument(file_path).paragraphs) // 50 + 1  # Примерный расчет
            elif file_extension == '.pdf':
                with open(file_path, 'rb') as f:
                    pages = len(PdfReader(f).pages)
            elif file_extension in ['.pptx', '.ppt']:
                pages = len(Presentation(file_path).slides)
            elif file_extension in ['.jpg', '.jpeg', '.png']:
                # Для изображений всегда 1 страница
                pages = 1
            elif file_extension in ['.xls', '.xlsx']:
                # Для Excel считаем количество листов
                pages = len(pd.ExcelFile(file_path).sheet_names)
            elif file_extension == '.txt':
                # Для текстовых файлов примерный расчет
                with open(file_path, 'r') as f:
                    lines = len(f.readlines())
                    pages = lines // 60 + 1  # ~60 строк на страницу
            print(is_duplex,pages)
        except Exception as e:
            print(f"Ошибка при обработке файла: {e}")
            pages = 1  # Если возникла ошибка, считаем как 1 страницу

        # Рассчитываем стоимость
        price_per_page =  terminal.price_per_page_bw
        total_price = pages * price_per_page

        # Создаем документ в базе данных
        document = DocModel.objects.create(
            title=uploaded_file.name,
            file=uploaded_file,
            pages=pages,
            terminal=terminal,
            is_color=is_color,
            is_duplex=is_duplex,
            price_document=total_price
        )
        print(total_price)
        return redirect("payment", doc_id=document.id)

    return render(request, 'upload.html', {
        "terminal": terminal,
        "supported_formats": ['.docx', '.pdf', '.pptx', '.jpg', '.jpeg', '.png', '.xls', '.xlsx', '.txt']
    })

def read_docx(file_path):
    with open(file_path, 'rb') as docx_file:
        content = docx_file.read()
    return content.decode('utf-8', errors='ignore')


def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()


def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content


def get_word_page_count(file_path):
    try:
        with zipfile.ZipFile(file_path, 'r') as docx:
            with docx.open('docProps/app.xml') as app_xml:
                dom = xml.dom.minidom.parse(app_xml)
                pages = dom.getElementsByTagName('Pages')[0].childNodes[0].nodeValue
                return int(pages)
    except Exception as e:
        return f"Error counting pages: {e}"



def count_pages_docx(file_path):
    doc = Document(file_path)
    page_breaks = sum([1 for para in doc.paragraphs if para.text == ''])
    return page_breaks + 1





def printer_list(request):
    printers = Printer.objects.all()
    return render(request, 'index.html', {'printers': printers})


def detail(request,pk):
    printers = Printer.objects.get(id=pk)
    return render(request, 'detail.html', {'printer': printers})


def process_printing(request, doc_id):
    """
    Обрабатывает печать документа и обновляет статус принтера
    """
    try:
        # Получаем документ и связанные устройства
        document = get_object_or_404(DocModel, id=doc_id)
        terminal = get_object_or_404(Terminal, id=document.terminal.id)
        printer = get_object_or_404(Printer, terminal=terminal.id)
        
        # Проверяем достаточно ли бумаги
        if printer.paper_count < document.pages:
            return render(request, "load.html", {
                'error': f"Недостаточно бумаги. Нужно {document.pages} листов, доступно {printer.paper_count}"
            })
        
        # Путь к файлу документа
        file_path = os.path.join(settings.MEDIA_ROOT, str(document.file))
        
        # Отправляем документ на печать через API
        # api_url = 'https://4c14-91-247-59-83.ngrok-free.app/api/upload/'
        api_url = terminal.ngrok_url
        files = {'file': open(file_path, 'rb')}
        data = {'title': document.title}
        
        response = requests.post(api_url, files=files, data=data)
        
        if response.status_code != 201:
            return render(request, "load.html", {
                'error': f"Ошибка печати. Код ошибки: {response.status_code}"
            })
        
        # Обновляем количество бумаги
        printer.paper_count -= document.pages
        printer.save()
        
        # Возвращаем успешный результат
        return render(request, "load.html", {
            'success': True,
            'document': document,
            'terminal': terminal,
            'printer': printer,
            'pages_printed': document.pages,
            'remaining_paper': printer.paper_count
        })
        
    except Exception as e:
        return render(request, "load.html", {
            'error': f"Произошла ошибка: {str(e)}"
        })
def load_document(request,doc_id):
    
    document = get_object_or_404(DocModel, id=doc_id)
    terminal2 = get_object_or_404(Terminal, id=document.terminal.id)
    printer2 = get_object_or_404(Printer, terminal=terminal2.id)
    printer2.paper_count -= document.pages
    printer2.save()
    return render(request, "load.html")
# 1️⃣ Список всех терминалов
# def terminal_list(request):
#     terminals = Terminal.objects.all()
#     data = [
#         {
#             "id": t.id,
#             "name": t.name,
#             "location": t.location,
#             "supports_color": t.supports_color,
#             "price_per_page_bw": t.price_per_page_bw,
#             "price_per_page_color": t.price_per_page_color,
#         }
#         for t in terminals
#     ]
#     return JsonResponse({"terminals": data})

# 2️⃣ Загрузка документа
# @csrf_exempt
# def upload_document(request, terminal_id):
#     if request.method == "POST":
#         terminal = get_object_or_404(Terminal, id=terminal_id)
#         file = request.FILES.get("file")
#         is_color = request.POST.get("is_color") == "true"

#         if not file:
#             return JsonResponse({"error": "Файл не загружен"}, status=400)

#         # Определяем количество страниц
#         extension = os.path.splitext(file.name)[1].lower()
#         pages = 1  # по умолчанию

#         try:
#             if extension == ".pdf":
#                 pages = get_pdf_page_count(file)
#             elif extension == ".pptx":
#                 pages = get_pptx_page_count(file)
#         except Exception:
#             return JsonResponse({"error": "Ошибка при обработке файла"}, status=400)

#         # Создаем объект документа
#         document = Document.objects.create(
#             title=file.name,
#             file=file,
#             pages=pages,
#             terminal=terminal,
#             is_color=is_color
#         )

#         return JsonResponse({"document_id": document.id, "pages": pages})

#     return JsonResponse({"error": "Метод не разрешен"}, status=405)

# Библиотека для работы с PDF

# def upload_document(request, pk):
#     terminal = get_object_or_404(Terminal, id=pk)  # Получаем терминал
#     price_per_page = terminal.price_per_page_bw  # Цена за 1 страницу (пример)
    
#     if request.method == "POST":
#         file = request.FILES.get("file")
#         is_color = request.POST.get("color_print") == "on"
        
#         # Определяем количество страниц (если PDF)
#         num_pages = 1  # По умолчанию 1 (если не PDF)
#         if file.name.endswith(".pdf"):
#             reader = PyPDF2.PdfReader(file)
#             num_pages = len(reader.pages)
#         print(num_pages)
#         total_price = num_pages * price_per_page  # Рассчитываем стоимость

#         # Сохраняем документ
#         document = Document.objects.create(
#             title=file.name,
#             file=file,
#             terminal=terminal,
#             is_color=is_color,
#             pages=num_pages,
#             price_document = total_price
            
#         )

#         # Перенаправляем на оплату, передавая ID документа
#         return redirect("payment", doc_id=document.id)

#     return render(request, "upload.html", {"terminal": terminal})
import os
import zipfile
from xml.dom import minidom
from django.http import JsonResponse
from django.shortcuts import render, get_object_or_404, redirect
from django.core.files.storage import FileSystemStorage
import mammoth
import pandas as pd
import PyPDF2
from docx import Document  # Импортируем для подсчета страниц в DOCX
from .models import Terminal, Document as DocModel


def upload_document(request, pk):
    terminal = get_object_or_404(Terminal, id=pk)  # Получаем терминал
    price_per_page = terminal.price_per_page_bw  # Цена за 1 страницу (пример)
    num_pages = 1  # По умолчанию 1 (если не PDF и не DOCX)

    if request.method == "POST" and request.FILES.get('file'):
        uploaded_file = request.FILES['file']
        is_color = request.POST.get("color_print") == "on"
        
        # Сохраняем файл
        fs = FileSystemStorage()
        filename = fs.save(uploaded_file.name, uploaded_file)
        file_path = fs.path(filename)

        # Определение типа файла
        file_extension = os.path.splitext(filename)[1].lower()

        if file_extension == '.docx':
            content = read_docx(file_path)
            num_pages = count_pages_docx(file_path)
        elif file_extension in ['.xls', '.xlsx']:
            content = read_excel(file_path)
        elif file_extension == '.txt':
            content = read_txt(file_path)
        elif file_extension == '.pdf':
            num_pages = count_pages_pdf(file_path)
        else:
            content = "Unsupported file format"

        # Рассчитываем стоимость
        print(type(num_pages),num_pages)
        from decimal import Decimal

        num_pages = int(num_pages)
        price_per_page = Decimal(price_per_page)
        total_price = num_pages * price_per_page

        print(total_price)
        # Сохраняем документ в базе данных
        document = DocModel.objects.create(
            title=uploaded_file.name,
            file=uploaded_file,
            terminal=terminal,
            is_color=is_color,
            pages=num_pages,
            price_document=total_price
        )

        # Перенаправляем на оплату, передавая ID документа
        return redirect("payment", doc_id=document.id)

    return render(request, "upload.html", {"terminal": terminal})


def read_docx(file_path):
    with open(file_path, 'rb') as docx_file:
        result = mammoth.extract_raw_text(docx_file)
        content = result.value
    return content


def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()


def read_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content


def count_pages_docx(file_path):
    """
    Подсчет страниц в DOCX-документе.
    """
    try:
        doc = Document(file_path)
        # Подсчет разрывов страниц
        page_breaks = sum([1 for para in doc.paragraphs if para.text == ''])
        return page_breaks + 1
    except Exception as e:
        return f"Error counting pages: {e}"


def count_pages_pdf(file_path):
    """
    Подсчет страниц в PDF-документе.
    """
    try:
        reader = PyPDF2.PdfReader(file_path)
        return len(reader.pages)
    except Exception as e:
        return f"Error counting pages: {e}"


def payment(request, doc_id):
    document = get_object_or_404(DocModel, id=doc_id)

    if request.method == "POST":
        # Тут будет логика обработки платежа
        return redirect("print_document", doc_id=document.id)

    return render(request, "payment.html", {"document": document})

# def upload_document(request, pk):
#     terminal = get_object_or_404(Terminal, id=pk)  # Получаем терминал по ID

#     if request.method == "POST":
#         file = request.FILES.get("file")
#         is_color = request.POST.get("color_print") == "on"

#         if file:
#             Document.objects.create(
#                 title=file.name,
#                 file=file,
#                 terminal=terminal,  # Используем терминал по ID
#                 is_color=is_color
#             )
#             return redirect("printer_list")  # Перенаправление после загрузки

#     return render(request, "upload.html", {"terminal": terminal})

# def upload_document(request):
#     terminals = Terminal.objects.all()

#     if request.method == "POST":
#         terminal_id = request.POST.get("terminal")
#         file = request.FILES.get("file")
#         is_color = request.POST.get("color_print") == "on"

#         if terminal_id and file:
#             terminal = Terminal.objects.get(id=terminal_id)
#             document = Document.objects.create(
#                 title=file.name,
#                 file=file,
#                 terminal=terminal,
#                 is_color=is_color
#             )
#             return redirect("success_page")  # Перенаправление после загрузки

#     return render(request, "upload.html", {"terminals": terminals})

# 3️⃣ Оплата документа
# @csrf_exempt
# def process_payment(request, document_id):
#     if request.method == "POST":
#         document = get_object_or_404(Document, id=document_id)
#         terminal = document.terminal

#         # Рассчитываем цену
#         price_per_page = terminal.price_per_page_color if document.is_color else terminal.price_per_page_bw
#         total_price = price_per_page * document.pages

#         # Создаем платеж
#         payment = Payment.objects.create(
#             terminal=terminal,
#             document=document,
#             amount=total_price,
#             status="pending",
#         )

#         # 🔗 Логика онлайн-оплаты (здесь должна быть интеграция с платежной системой)
#         payment.status = "completed"  # Симуляция успешного платежа
#         payment.save()

#         return JsonResponse({"payment_id": payment.id, "status": payment.status})

#     return JsonResponse({"error": "Метод не разрешен"}, status=405)


# 4️⃣ Запуск печати
@csrf_exempt
def print_document(request, document_id):
    if request.method == "POST":
        document = get_object_or_404(Document, id=document_id)
        printer = Printer.objects.filter(terminal=document.terminal).first()

        if not printer:
            return JsonResponse({"error": "Принтер не найден"}, status=404)

        if printer.paper_count < document.pages:
            return JsonResponse({"error": "Недостаточно бумаги"}, status=400)

        # 🔥 Отправляем команду на печать
        send_to_printer(document.file.path)

        # Уменьшаем количество бумаги
        printer.paper_count -= document.pages
        printer.save()

        return JsonResponse({"status": "Печать начата"})

    return JsonResponse({"error": "Метод не разрешен"}, status=405)


# 📌 Функции для обработки файлов
def get_pdf_page_count(file):
    file.seek(0)
    pdf = PdfReader(file)
    return len(pdf.pages)

def get_pptx_page_count(file):
    file.seek(0)
    ppt = Presentation(file)
    return len(ppt.slides)

# 📌 Функция отправки на печать (нужно заменить на реальную)
def send_to_printer(file_path):
    os.system(f"lp {file_path}")  # Команда для Linux


# from django.http import HttpResponse
# import requests

#  # Примерные данные, которые отправляешь

# from django.middleware.csrf import get_token

# # Проверка ответа



# def index(request):
#     return render(request, 'main.html')

# from django.shortcuts import render, redirect
# from .forms import DocumentForm
# from .models import Document, PrinterStatus
# # headers = {'X-CSRFToken': csrf_token}
# def upload_document(request):
#     printer_status = PrinterStatus.objects.first()  # Получаем состояние аппарата

#     if request.method == 'POST':
#         form = DocumentForm(request.POST, request.FILES)
#         if form.is_valid():
#             document = form.save(commit=False)
#             document.pages = form.cleaned_data['pages']
#             document.title = form.cleaned_data['title']
#             document.file = form.cleaned_data['file']
#             document.printer = form.cleaned_data['printer']
#             document.save()
#             file_document = document.file.path

#             url = 'https://4c14-91-247-59-83.ngrok-free.app/api/upload/'

#                 # Данные для отправки
#                 # files = {'file': open('document.pdf', 'rb')}
#             files = {'file':  open( file_document, 'rb') }
#             data = {'title': document.title}

#                 # Отправка POST-запроса

#             response = requests.post(url, files=files, data=data)
#             # Проверяем, хватает ли бумаги
#             if printer_status and document.pages > printer_status.paper_amount:
#                 message = f"Не хватает бумаги для печати документа ({document.pages} страниц). Осталось {printer_status.paper_amount} листов."
#             else:
#                 message = f"Документ успешно загружен! Страниц: {document.pages}."
#                 # URL сервера на ноутбуке


#                 if response.status_code == 201:
#                     message += 'Документ успешно отправлен на печать.'
#                 else:
#                     message += f'Ошибка при отправке документа: {response.status_code}'
#                 # Обновляем количество бумаги после печати
#                 if printer_status:
#                     printer_status.paper_amount -= document.pages
#                     printer_status.save()

#             return render(request, 'upload.html', {'form': form, 'message': message})
#     else:
#         form = DocumentForm()

#     return render(request, 'upload.html', {'form': form, 'printer_status': printer_status})

# from django.core.files.storage import default_storage

# def server(request):

#     # Получаем документ с id=50
#     doc = Document.objects.get(id=50)
#     # URL сервера на ноутбуке
#     url = 'https://4c14-91-247-59-83.ngrok-free.app/api/upload/'
#     # Данные для отправки
#     data = {'title': 'Мой документ'}
#     # Получаем полный путь к файлу
#     file_path = doc.file.path  # Полный путь к файлу
#     # Открываем файл для отправки
#     with open(file_path, 'rb') as f:
#         files = {'file': f}
#         # Отправляем POST-запрос на сервер с файлом и данными
#         response = requests.post(url,  files=files, data=data)
#     # Обработка ответа
#     message = '___'
#     if response.status_code == 201:
#         message += 'Документ успешно отправлен на печать.'
#     else:
#         message += f'{file_path} Ошибка при отправке документа: {response.status_code}'

#     # except Document.DoesNotExist:
#     #     message = 'Документ не найден.'

#     # except FileNotFoundError:
#     #     message = 'Файл документа не найден.'

#     # except Exception as e:
#     #     message = f'Произошла ошибка: {str(e)}'

#     return render(request, 'go.html', {'message': message})
#     # doc = Document.objects.get(id=50)
#     # url = 'http://192.168.0.104:8000/api/upload/'

#     #             # Данные для отправки
#     #             # files = {'file': open('document.pdf', 'rb')}

#     # data = {'title': 'Мой документ',}
#     # file_path = doc.file.path  # Полный путь к файлу

#     # # Открываем файл для отправки
#     # # with default_storage.open(file_path, 'rb') as file:
#     # files = {'file': open('test.pdf', 'rb')}
#     # response = requests.post(url, files=files, data=data)
#     # message = '___'
#     # if response.status_code == 201:
#     #     message += 'Документ успешно отправлен на печать.'
#     # else:
#     #     message += f'{file_path}Ошибка при отправке документа: {response.status_code}-'

#     # return render(request, 'go.html', {'message': message})

# def document_list(request):
#     documents = Document.objects.all()
#     return render(request, 'document_list.html', {'documents': documents})

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from .models import Terminal

@csrf_exempt
def update_terminal(request, terminal_id):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            ngrok_url = data.get("ngrok_url")

            terminal = Terminal.objects.get(id=terminal_id)
            terminal.ngrok_url = ngrok_url
            terminal.save()

            return JsonResponse({"status": "ok", "ngrok_url": ngrok_url})
        except Exception as e:
            return JsonResponse({"status": "error", "message": str(e)})
    return JsonResponse({"status": "invalid request"}, status=400)


from django.shortcuts import render, get_object_or_404
from .models import TemplateDocumentGroup

def document_groups(request, group_id=None):
    # Если `group_id` передан, находим текущую группу или показываем корневые группы
    if group_id:
        current_group = get_object_or_404(TemplateDocumentGroup, id=group_id)
        subgroups = current_group.subgroups.all()
        documents = current_group.documents.all()
    else:
        current_group = None
        subgroups = TemplateDocumentGroup.objects.filter(parent_group__isnull=True)
        documents = []

    return render(request, "shablon/document_tree.html", {
        "current_group": current_group,
        "subgroups": subgroups,
        "documents": documents,
    })

from django.shortcuts import render
from .models import TemplateDocumentGroup

def document_tree(request):
    groups = TemplateDocumentGroup.objects.filter(parent_group__isnull=True)  # Только корневые группы
    return render(request, "shablon/list.html", {"groups": groups})

def print_document(request, document_id):
    document = get_object_or_404(TemplateDocument, id=document_id)

    # Здесь вы можете реализовать выбор принтера и оплату.
    if request.method == "POST":
        selected_printer = request.POST.get("printer")
        return render(request, "documents/payment.html", {
            "document": document,
            "printer": selected_printer,
        })

    return render(request, "documents/select_printer.html", {
        "document": document,
        "printers": ["Принтер №1", "Принтер №2", "Принтер №3"],  # Пример списка принтеров.
    })


# import requests
# from .models import Terminal

# def check_terminals():
#     terminals = Terminal.objects.all()
    
#     for terminal in terminals:
#         if not terminal.ngrok_url:
#             terminal.is_active = False
#             terminal.save()
#             continue

#         try:
#             response = requests.get(f"{terminal.ngrok_url}/status", timeout=5)
#             if response.status_code == 200:
#                 terminal.is_active = True
#             else:
#                 terminal.is_active = False
#         except requests.RequestException:
#             terminal.is_active = False

#         terminal.save()


def template_list(request):
    templates = TemplateDocument.objects.all()
    return render(request, "shablon/list.html", {"templates": templates})

#@login_required
def select_terminal(request, document_id):
    document = get_object_or_404(TemplateDocument, id=document_id)
    terminals = Terminal.objects.filter(is_active=True)
    return render(request, "shablon/select_terminal.html", {"document": document, "terminals": terminals})

#@login_required
def pay_and_print(request, document_id, terminal_id):
    document = get_object_or_404(TemplateDocument, id=document_id)
    terminal = get_object_or_404(Terminal, id=terminal_id)

    # Создаем задачу печати
    job = PrintJob.objects.create(user=request.user, document=document, terminal=terminal, status="paid")

    # Отправляем запрос на терминал для печати
    api_url = f"{terminal.ngrok_url}/print/"
    data = {
        "document_url": request.build_absolute_uri(document.file.url),
        "pages": 1,  # Можно добавить анализ количества страниц
        "user_id": request.user.id,
    }
    try:
        response = requests.post(api_url, json=data, timeout=5)
        if response.status_code == 200:
            job.status = "printed"
            job.save()
    except requests.exceptions.RequestException:
        pass
    return redirect('print_document', doc_id=document.id)
    # return render(request, "shablon/print_success.html", {"document": document, "terminal": terminal})